from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Define a helper function to add bullet slides
def add_bullet_slide(title_text, bullet_points, font_size=24):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = title_text

    # Clear default text
    content.text = ""
    tf = content.text_frame
    tf.word_wrap = True

    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(font_size)
        p.font.name = "Calibri"
        p.font.color.rgb = RGBColor(0, 51, 102)  # Navy blue

# --- Slide 1: Cover ---
slide_layout = prs.slide_layouts[0]  # Title Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Prepositions: Between, Among, Under"
subtitle.text = "Usage, Non-Usage, and Types"

# --- Slide 2: Index ---
add_bullet_slide("Index", [
    "Types of Prepositions",
    "Between",
    "Among",
    "Under",
    "Summary"
])

# --- Slide 3: Types of Prepositions ---
add_bullet_slide("Types of Prepositions", [
    "Preposition of Time",
    "Preposition of Place",
    "Preposition of Agent",
    "Preposition of Manner",
    "Preposition of Instrument",
    "Preposition of Cause/Purpose",
    "Preposition of Direction/Movement",
    "Preposition of Measure (extended lists)",
    "Focus today: Place/Position → Between, Among, Under"
])

# --- Slide 4: Between ---
add_bullet_slide("Between", [
    "Type: Preposition of Place/Position",
    "Used for two distinct items or people",
    "Examples: 'The ball is between the two boxes'",
    "Examples: 'She sat between her parents'",
    "Not used for groups larger than two",
    "Incorrect: 'She was popular between classmates'",
    "Correct: 'She was popular among classmates'"
])

# --- Slide 5: Among ---
add_bullet_slide("Among", [
    "Type: Preposition of Place/Position",
    "Used for more than two items in a group",
    "Examples: 'The deer was hiding among the trees'",
    "Examples: 'She was popular among her classmates'",
    "Not used for only two items",
    "Incorrect: 'The ball is among two boxes'",
    "Correct: 'The ball is between two boxes'"
])

# --- Slide 6: Under ---
add_bullet_slide("Under", [
    "Type: Preposition of Place/Position",
    "Refers to something below or covered",
    "Examples: 'The cat is under the table'",
    "Examples: 'The book is under the pillow'",
    "Not used for 'inside' or 'between'",
    "Incorrect: 'The pen is under the bag’s pocket'",
    "Correct: 'The pen is inside the bag’s pocket'"
])

# --- Slide 7: Summary ---
add_bullet_slide("Summary", [
    "Between → Place/Position (two items)",
    "Among → Place/Position (group/more than two)",
    "Under → Place/Position (below/covered)",
    "Key Tip: All three are prepositions of place/position",
    "Difference: Number (two vs. group) and Location (below vs. inside)"
])

# Save the presentation
prs.save("Prepositions_PPT_Designed.pptx")
print("Presentation created successfully: Prepositions_PPT_Designed.pptx")
